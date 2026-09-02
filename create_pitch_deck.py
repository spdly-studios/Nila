from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = 'attendly_pitch_deck.pptx'
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BG = RGBColor(247,250,248); INK = RGBColor(23,33,43); MUTED = RGBColor(113,128,139); GREEN = RGBColor(25,115,77); MINT = RGBColor(216,240,229); CORAL = RGBColor(226,120,91); NAVY = RGBColor(38,55,70); WHITE = RGBColor(255,255,255)

def box(slide, x,y,w,h, fill=None, line=None, radius=False):
    sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)); sh.fill.solid(); sh.fill.fore_color.rgb=fill or BG; sh.line.color.rgb=line or (fill or BG); return sh
def text(slide, s, x,y,w,h, size=18, color=INK, bold=False, font='Aptos', align=PP_ALIGN.LEFT):
    if s == '2' and isinstance(h, RGBColor): x,y,w,h,size,color,bold = 9.15,2.2,.4,.4,22,h,True
    tb=slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf=tb.text_frame; tf.clear(); tf.word_wrap=True; tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0; p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=s; r.font.name=font; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color; return tb
def base(kicker, num):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); sl.background.fill.solid(); sl.background.fill.fore_color.rgb=BG; text(sl,kicker.upper(),.65,.42,5,.25,10,GREEN,True); text(sl,f'{num:02d}',12.1,.42,.5,.25,10,MUTED,True,align=PP_ALIGN.RIGHT); return sl
def bullet(slide, s, x,y,w, size=18, color=INK):
    text(slide,'—',x,y,.22,.3,size,CORAL,True); text(slide,s,x+.35,y,w-.35,.45,size,color)

sl=prs.slides.add_slide(prs.slide_layouts[6]); sl.background.fill.solid(); sl.background.fill.fore_color.rgb=NAVY; box(sl,.62,.62,.12,6.25,CORAL); text(sl,'attendly',1.05,.78,5,.5,22,WHITE,True,'Aptos Display'); text(sl,'Attendance that explains\nitself.',1.05,1.65,9,1.55,42,WHITE,True,'Aptos Display'); text(sl,'A teacher should not have to guess why a student is absent.',1.08,3.55,7,.5,18,RGBColor(216,240,229)); box(sl,8.9,1.55,3.5,3.5,MINT,line=MINT,radius=True); text(sl,'NILA',9.35,2.15,2.5,.35,12,GREEN,True); text(sl,'A calmer\nfollow-up loop',9.35,2.65,2.7,1.2,27,NAVY,True,'Aptos Display'); text(sl,'SpDly Studios  ·  Shivaprasad V',1.08,6.65,6,.3,11,RGBColor(184,204,200));

sl=base('The gap',2); text(sl,'Attendance tells you who is missing.\nIt rarely tells you why.',.7,1.1,8.6,1.2,32,INK,True,'Aptos Display'); bullet(sl,'A mark in a register is a signal—not an explanation.',.75,2.75,5.5); bullet(sl,'A phone call may capture the truth, but the insight gets buried in memory, notes, or a transcript.',.75,3.4,6.5); bullet(sl,'Patterns across weeks, students, and weekdays remain invisible.',.75,4.25,6.5); box(sl,8.5,1.25,3.8,4.7,WHITE,line=RGBColor(220,230,226),radius=True); text(sl,'ABSENT',9.0,1.9,2.8,.3,12,CORAL,True); text(sl,'?',9.55,2.45,1.5,1.2,72,CORAL,True,'Aptos Display',PP_ALIGN.CENTER); text(sl,'The missing layer\nis context.',9.0,4.25,2.8,.7,20,NAVY,True,'Aptos Display',PP_ALIGN.CENTER)

sl=base('The product',3); text(sl,'Attendly closes the loop\nbetween absence and action.',.7,1.05,8.6,1.25,32,INK,True,'Aptos Display'); text(sl,'A single classroom workspace that records attendance, reaches parents, captures what they say, and surfaces the patterns a busy teacher cannot manually calculate.',.75,2.7,7.3,1.0,18,MUTED); box(sl,8.65,1.15,3.7,4.8,NAVY,line=NAVY,radius=True); text(sl,'1',9.15,1.7,.4,.4,22,RGBColor(244,189,124),True); text(sl,'Mark',9.8,1.72,2,.3,16,WHITE,True); text(sl,'2',2.22,.4,22,RGBColor(244,189,124),True); text(sl,'Reach',9.8,2.25,2,.3,16,WHITE,True); text(sl,'3',9.15,2.8,.4,.4,22,RGBColor(244,189,124),True); text(sl,'Understand',9.8,2.82,2,.3,16,WHITE,True); text(sl,'4',9.15,3.9,.4,.4,22,RGBColor(244,189,124),True); text(sl,'Act',9.8,3.92,2,.3,16,WHITE,True); text(sl,'Less chasing.\nMore context.',9.15,4.8,2.5,.7,20,MINT,True,'Aptos Display')

sl=base('The experience',4); text(sl,'One workflow. Four moments.',.7,1.0,8,0.6,31,INK,True,'Aptos Display'); steps=[('01','Register','Teacher marks attendance by class.'),('02','Connect','Nila calls the parent in Tamil first.'),('03','Capture','Webhook stores transcript, summary, disposition, and recording metadata.'),('04','Learn','AI extracts the parent’s main reason and tags; analytics reveal recurrence.')]; x=.8
for n,t,d in steps:
    box(sl,x,2.05,2.75,3.3,WHITE,line=RGBColor(220,230,226),radius=True); text(sl,n,x+.25,2.35,.5,.3,12,CORAL,True); text(sl,t,x+.25,2.9,2.2,.4,21,NAVY,True,'Aptos Display'); text(sl,d,x+.25,3.7,2.2,1.0,14,MUTED); x+=3.05

sl=base('The intelligence layer',5); text(sl,'The reason is not always\nin the first sentence.',.7,1.0,7.6,1.0,31,INK,True,'Aptos Display'); text(sl,'Attendly keeps the original evidence, then adds a readable layer for the teacher.',.75,2.25,6.6,.5,17,MUTED); box(sl,.8,3.25,5.5,2.4,WHITE,line=RGBColor(220,230,226),radius=True); text(sl,'Parent-only summary',1.15,3.65,2.6,.3,12,GREEN,True); text(sl,'“Transport issue affected\nMeera’s attendance today.”',1.15,4.2,4.3,.9,23,NAVY,True,'Aptos Display'); box(sl,6.75,3.25,5.5,2.4,MINT,line=MINT,radius=True); text(sl,'Signals for action',7.1,3.65,2.6,.3,12,GREEN,True); text(sl,'transport   ·   recurring   ·   Wednesday',7.1,4.2,4.5,.8,20,NAVY,True,'Aptos Display');

sl=base('Teacher value',6); text(sl,'Useful before, during,\nand after the school day.',.7,1.0,8,1.0,31,INK,True,'Aptos Display'); vals=[('Before class','See who needs follow-up without rebuilding a list.'),('During the day','Let Nila handle routine conversations with a consistent, respectful flow.'),('After the call','Read the main reason—not a wall of transcript—and decide what needs attention.'),('Over time','Spot repeated weekday, transport, health, or group patterns early.')]; y=2.35
for title,desc in vals: text(sl,title,.9,y,2.0,.3,15,CORAL,True); text(sl,desc,3.0,y,8.6,.4,17,INK); y+=.8

sl=base('Trust by design',7); text(sl,'Automation with an evidence trail.',.7,1.0,8,0.6,31,INK,True,'Aptos Display'); bullet(sl,'Original SnapServe transcript stays available.',.8,2.1,6.3); bullet(sl,'AI summary uses parent/caller speech—not agent instructions or memory.',.8,2.8,7); bullet(sl,'Tamil-first conversations switch only when the parent asks.',.8,3.5,6.5); bullet(sl,'Webhook, sync, analysis, and errors are logged for debugging.',.8,4.2,7); box(sl,9.1,1.7,2.5,3.4,NAVY,line=NAVY,radius=True); text(sl,'EVIDENCE',9.45,2.2,1.8,.3,11,RGBColor(244,189,124),True); text(sl,'source\n→\nreason\n→\naction',9.45,2.8,1.8,1.8,23,MINT,True,'Aptos Display',PP_ALIGN.CENTER)

sl=base('Live proof',8); text(sl,'A working product, not a promise.',.7,1.0,8,0.6,31,INK,True,'Aptos Display'); box(sl,.8,2.0,3.5,2.9,WHITE,line=RGBColor(220,230,226),radius=True); text(sl,'PROJECT',1.15,2.45,1.5,.25,10,GREEN,True); text(sl,'attendly',1.15,2.9,2.2,.45,25,NAVY,True,'Aptos Display'); text(sl,'AI voice agent: Nila',1.15,3.75,2.3,.3,15,MUTED); text(sl,'Team: SpDly Studios',1.15,4.2,2.3,.3,15,MUTED); box(sl,4.7,2.0,7.6,2.9,MINT,line=MINT,radius=True); text(sl,'DEPLOYED',5.1,2.45,1.5,.25,10,GREEN,True); text(sl,'Open the live call history',5.1,2.95,4.5,.4,22,NAVY,True,'Aptos Display'); text(sl,'nila-v80b.onrender.com/history.html',5.1,3.8,6.4,.35,16,GREEN,True); text(sl,'Built for teachers. Designed for clarity.',5.1,4.45,5.5,.3,15,NAVY)

sl=prs.slides.add_slide(prs.slide_layouts[6]); sl.background.fill.solid(); sl.background.fill.fore_color.rgb=NAVY; text(sl,'The ask',.9,.85,3,.3,11,RGBColor(244,189,124),True); text(sl,'Make every absence\nunderstandable.',.9,1.55,8,1.2,38,WHITE,True,'Aptos Display'); text(sl,'Attendly gives teachers the missing context to respond earlier, more thoughtfully, and with less administrative work.',.95,3.25,7.3,.8,19,MINT); text(sl,'SpDly Studios',.95,6.25,3,.3,14,WHITE,True); text(sl,'Shivaprasad V  ·  Nila  ·  attendly',.95,6.65,5,.3,11,RGBColor(184,204,200)); text(sl,'Live demo →',10.2,6.35,2.2,.3,14,RGBColor(244,189,124),True,align=PP_ALIGN.RIGHT)

prs.save(OUT)
print(OUT)
